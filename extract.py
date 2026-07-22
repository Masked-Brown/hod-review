#!/usr/bin/env python3
"""extract.py — deterministic PowerPoint reader (manifest M3).

Input a `.pptx`; output a **slide manifest**: the verbatim text of each slide,
numbered in document order. Stdlib plus python-pptx only.

Why deterministic. Slide anchoring and verbatim-quote checking are exactly the kind
of matching that must never run on model diligence (`brainwave.md`, "Why a binary
input is the wedge"). This reader does the anchoring once, mechanically, so that
`check.py` can verify every `QUOTE` in a critique against the real slide text
(`reference/finding-schema.md` §3, the fabrication check).

The manifest format (consumed by `check.py`'s `parse_manifest`):

    # Slide manifest
    # source: <basename>
    # slides: <N>
    # generated-by: extract.py ...

    ## Slide 1
    <verbatim text of slide 1, one shape-chunk per line>

    ## Slide 2
    <verbatim text of slide 2>
    ...

`check.py` locates a slide by the exact header line `## Slide <n>` and reads every
line up to the next such header as that slide's text. A slide with no text still
gets a `## Slide <n>` header (with an empty body) so slide numbering never drifts.

Usage:
    python extract.py deck.pptx            # manifest to stdout
    python extract.py deck.pptx -o out.md  # manifest to a file

Text coverage: title and body placeholders, text boxes, tables (cell by cell), and
text nested inside grouped shapes. Charts, images and speaker notes carry no
on-slide text and are not extracted — a `QUOTE` must be wording a student could read
off the slide.
"""

from __future__ import annotations

import argparse
import os
import sys

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:  # pragma: no cover - environment guard
    sys.stderr.write(
        "extract.py requires python-pptx (pip install python-pptx). "
        "Stdlib + python-pptx only.\n"
    )
    raise


def _shape_texts(shape) -> list[str]:
    """Verbatim text chunks from one shape, recursing into groups and tables.

    Returns a list of strings in document order. Each chunk is emitted verbatim;
    no tidying, no case changes, no punctuation stripping — the fabrication check
    downstream depends on this being the real words in the real order.
    """
    texts: list[str] = []

    shape_type = getattr(shape, "shape_type", None)
    if shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            texts.extend(_shape_texts(sub))
        return texts

    if getattr(shape, "has_table", False):
        table = shape.table
        for row in table.rows:
            cells = [cell.text for cell in row.cells]
            texts.append("\t".join(cells))
        return texts

    if getattr(shape, "has_text_frame", False):
        # text_frame.text joins paragraphs with "\n" and runs within a paragraph
        # with no separator — i.e. exactly what a reader sees on the slide.
        frame_text = shape.text_frame.text
        if frame_text:
            texts.append(frame_text)
        return texts

    return texts  # picture, chart, connector, etc. — no on-slide text


def _slide_text_lines(slide) -> list[str]:
    """The verbatim text of one slide as a list of lines (may be empty)."""
    chunks: list[str] = []
    for shape in slide.shapes:
        chunks.extend(_shape_texts(shape))
    joined = "\n".join(c for c in chunks if c)
    if joined == "":
        return []
    return [line.rstrip() for line in joined.split("\n")]


def extract_manifest(pptx_path: str) -> str:
    """Read a .pptx and return the slide-manifest text (see module docstring)."""
    prs = Presentation(pptx_path)
    slides = list(prs.slides)

    out: list[str] = [
        "# Slide manifest",
        f"# source: {os.path.basename(pptx_path)}",
        f"# slides: {len(slides)}",
        "# generated-by: extract.py (python-pptx) — deterministic, verbatim slide text",
        "",
    ]
    for index, slide in enumerate(slides, start=1):
        out.append(f"## Slide {index}")
        out.extend(_slide_text_lines(slide))
        out.append("")  # one blank line separates slides
    return "\n".join(out).rstrip() + "\n"


def main(argv: list[str]) -> int:
    for stream in (sys.stdout, sys.stderr):
        try:
            stream.reconfigure(encoding="utf-8", errors="replace")
        except (AttributeError, ValueError):  # pragma: no cover
            pass
    parser = argparse.ArgumentParser(
        prog="extract.py",
        description="Deterministic .pptx -> slide manifest (verbatim text per slide).",
    )
    parser.add_argument("pptx", help="path to the PowerPoint (.pptx) file")
    parser.add_argument(
        "-o", "--output", help="write the manifest here (default: stdout)"
    )
    args = parser.parse_args(argv)

    if not os.path.exists(args.pptx):
        sys.stderr.write(f"extract.py: no such file: {args.pptx}\n")
        return 2

    try:
        manifest = extract_manifest(args.pptx)
    except Exception as exc:  # noqa: BLE001 - report cleanly, do not traceback at users
        sys.stderr.write(f"extract.py: could not read {args.pptx}: {exc}\n")
        return 2

    if args.output:
        with open(args.output, "w", encoding="utf-8", newline="\n") as handle:
            handle.write(manifest)
        sys.stderr.write(f"extract.py: wrote {args.output}\n")
    else:
        sys.stdout.write(manifest)
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
