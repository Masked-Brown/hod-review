#!/usr/bin/env python3
"""tests/fixture_deck.py — the canonical gate-test deck (manifest M3).

This is a **test fixture for the enforcement gate**, not the constructed validation
run. That run — a full deck, its answer key, and the CONSTRUCTED training rows — is
manifest M4's job and lives under `runs/` (`plan.md` §M4). This module exists only so
`tests/verify.py` can exercise the real pipeline end to end:

    build a tiny .pptx  ->  extract.py  ->  slide manifest  ->  check.py

The deck is deliberately small (six slides) and its text is fixed here, so the clean
critique in `tests/cases/` and the broken critiques in `tests/negative/` can quote it
verbatim and the fabrication check has a real, known ground truth to match against.

The slides borrow the validated topic (AQA 3.2.1 Cell structure) purely so the
fixtures can cite real handles from `reference/`. Three slides carry a seeded issue a
finding can legitimately anchor to; slide 6 is clean and draws no finding. The seeded
issues are gate-fixture scaffolding, not pedagogy under review.

Requires python-pptx (already a build dependency of extract.py).
"""

from __future__ import annotations

# Each slide: (title, [body lines]). Text is verbatim ground truth for the fixtures.
SLIDES: list[tuple[str, list[str]]] = [
    ("Cell Structure", [
        "AQA A-level Biology — 3.2",
        "Starter: copy the title and today's learning objectives",
    ]),
    ("Learning objectives", [
        "Describe the ultrastructure of eukaryotic cells",
        "Explain how the electron microscope reveals internal detail",
    ]),
    # Slide 3 — vocabulary trap (EI-1): "produce energy" is an explicit AQA reject.
    ("Mitochondria", [
        "Mitochondria are the powerhouse of the cell and produce energy for respiration.",
    ]),
    # Slide 4 — misconception (EI-5): attributes the optical limit to magnification.
    ("Seeing smaller structures", [
        "To see smaller organelles, increase the magnification.",
    ]),
    # Slide 5 — equation straight to independent work, no worked model (Rosenshine 4).
    ("Magnification calculation", [
        "Magnification = image size / actual size",
        "Now try the questions on the worksheet.",
    ]),
    # Slide 6 — clean; a correct plenary that should draw no finding.
    ("Plenary", [
        "Any questions before we finish?",
    ]),
]

FIXTURE_BASENAME = "gate-fixture-cell-structure.pptx"


def build(path: str) -> str:
    """Build the fixture deck at `path` and return `path`."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[6]  # fully blank: we place our own text boxes
    for title, body in SLIDES:
        slide = prs.slides.add_slide(blank)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
        title_box.text_frame.text = title
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(4))
        frame = body_box.text_frame
        for line_no, line in enumerate(body):
            paragraph = frame.paragraphs[0] if line_no == 0 else frame.add_paragraph()
            paragraph.text = line
    prs.save(path)
    return path
