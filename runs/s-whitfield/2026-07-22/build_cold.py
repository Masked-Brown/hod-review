#!/usr/bin/env python3
"""build_cold.py — build the CONSTRUCTED cold-run deck for the M4 run.

A SECOND, short constructed deck the editor has not seen, on a DIFFERENT sub-topic
(cell fractionation — the very area the main lesson.pptx skipped). Its purpose is to
show the editor generalises to unseen slides and unseen reference anchors (spec §5,
EI-6, EI-9, exam Q1), rather than being tuned to the first deck.

CONSTRUCTED — not a real lesson. Two slides seeded with a named flaw, two clean.
Text is ASCII so every QUOTE matches verbatim under check.py. Requires python-pptx.

Usage:  python build_cold.py            # writes cold-lesson.pptx next to this script
"""

from __future__ import annotations

import os

SLIDES: list[tuple[str, list[str]]] = [
    # Slide 1 — CLEAN — title.
    ("Cell fractionation", [
        "How we separate organelles to study them",
        "AQA A-level Biology 3.2.1",
    ]),
    # Slide 2 — SEED C1 — homogenisation solution conditions compressed to one vague line (EI-9).
    ("Preparing the sample", [
        "Blend the tissue to break open the cells, then filter it.",
        "Keep it in a suitable solution so the organelles are not damaged.",
    ]),
    # Slide 3 — SEED C2 — high-speed-first, no slow-first principle; conflates the technique (EI-6).
    ("Separating the organelles", [
        "Spin the filtered sample at high speed to separate the organelles by size.",
        "The heaviest organelles collect at the bottom.",
    ]),
    # Slide 4 — CLEAN — a correct closed-book retrieval plenary.
    ("Plenary", [
        "Close your notes.",
        "In order, list the three conditions of the solution and say why each one matters.",
    ]),
]

BASENAME = "cold-lesson.pptx"


def build(path: str) -> str:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[6]
    for title, body in SLIDES:
        slide = prs.slides.add_slide(blank)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
        title_box.text_frame.text = title
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(4.5))
        frame = body_box.text_frame
        for line_no, line in enumerate(body):
            paragraph = frame.paragraphs[0] if line_no == 0 else frame.add_paragraph()
            paragraph.text = line
    prs.save(path)
    return path


if __name__ == "__main__":
    here = os.path.dirname(os.path.abspath(__file__))
    out = os.path.join(here, BASENAME)
    build(out)
    print(f"built {out} ({len(SLIDES)} slides)")
