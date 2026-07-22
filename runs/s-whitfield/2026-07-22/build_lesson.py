#!/usr/bin/env python3
"""build_lesson.py — build the CONSTRUCTED validation deck for the M4 run.

This script deterministically builds `lesson.pptx`, the 12-slide **constructed**
demonstration deck reviewed in the s-whitfield / 2026-07-22 run. It is committed
alongside the .pptx so the deck is reproducible and its exact wording is auditable
(the fabrication check in check.py depends on the slide text being known ground truth).

CONSTRUCTED — not a real lesson. Some slides are deliberately clean; others are seeded
with a named flaw that maps to a specific entry in editor/reference/. The full map of
seeded flaw -> reference anchor is `expected-findings.md` (the answer key).

Text is authored in ASCII (straight quotes, hyphens) so every QUOTE in critique.md
matches verbatim under check.py's normalisation. Requires python-pptx.

Usage:  python build_lesson.py            # writes lesson.pptx next to this script
"""

from __future__ import annotations

import os

# Each slide: (title, [body paragraphs]). This text is the verbatim ground truth
# that slide-manifest.md, expected-findings.md, and critique.md all refer to.
SLIDES: list[tuple[str, list[str]]] = [
    # Slide 1 — CLEAN — title.
    ("AQA A-level Biology - 3.2 Cell structure", [
        "Lesson 1 of the Cells unit",
        "Ultrastructure, prokaryotes, viruses and microscopy",
        "Year 12",
    ]),
    # Slide 2 — SEED F1 — starter opens on new content, no retrieval (Rosenshine 1).
    ("Starter", [
        "Copy the organelle summary table from the sheet into your notes.",
        "We will use it through today's lesson.",
    ]),
    # Slide 3 — SEED F2 — lesson objectives skip cell fractionation (a 3.2.1 spec area).
    ("Today's lesson", [
        "By the end you will be able to:",
        "describe the ultrastructure of eukaryotic cells",
        "describe the structure of prokaryotic cells",
        "explain why viruses are acellular",
        "use a light microscope and calculate magnification",
    ]),
    # Slide 4 — CLEAN — correct overview of eukaryotic cells.
    ("The eukaryotic cell", [
        "Eukaryotic cells have membrane-bound organelles, each carrying out a specific function.",
        "Compartmentalisation lets incompatible reactions happen at the same time in one cell.",
    ]),
    # Slide 5 — SEED F3 — "produce energy" is an explicit AQA reject (EI-1).
    ("Mitochondria", [
        "The powerhouse of the cell.",
        "Mitochondria produce energy for the cell's reactions.",
        "Site of aerobic respiration.",
    ]),
    # Slide 6 — SEED F4 — "controls all the cell's activities" is GCSE-level, zero credit (EI-7).
    ("The nucleus", [
        "The control centre of the cell.",
        "The nucleus controls all the cell's activities.",
        "Surrounded by a nuclear envelope with pores.",
    ]),
    # Slide 7 — CLEAN — correct rER / Golgi, full terms, distinct functions.
    ("Rough ER and Golgi apparatus", [
        "Ribosomes on the rough endoplasmic reticulum synthesise proteins that are folded for secretion.",
        "The Golgi apparatus modifies these proteins and packages them into vesicles.",
    ]),
    # Slide 8 — SEED F5 — universal and non-universal prokaryote features mixed, unflagged (EI-2, the list rule).
    ("Prokaryotic cells", [
        "Features of prokaryotic cells:",
        "murein cell wall, cell-surface membrane, 70S ribosomes, circular DNA, capsule, plasmids, flagella",
    ]),
    # Slide 9 — SEED F6a + F6b — acellular/non-living fused (EI-3); "genetic information" a reject term.
    ("Viruses", [
        "Viruses are acellular and non-living because they cannot replicate outside a host cell.",
        "A virus is genetic information surrounded by a protein capsid.",
    ]),
    # Slide 10 — SEED F7 (CRITICAL) — magnification written for resolution (EI-5).
    ("Seeing smaller structures", [
        "The light microscope cannot show small organelles.",
        "To see them, increase the magnification for more detail.",
    ]),
    # Slide 11 — SEED F8 — equation then straight to independent work, no worked model (Rosenshine 4 / exam Q7).
    ("Magnification calculation", [
        "magnification = image size / actual size",
        "Now work through questions 1-6 on the worksheet.",
    ]),
    # Slide 12 — CLEAN — a genuine closed-book retrieval plenary.
    ("Plenary", [
        "Close your notes.",
        "On your whiteboard: name three organelles and give the precise function of each.",
        "Swap and check against the mark scheme.",
    ]),
]

BASENAME = "lesson.pptx"


def build(path: str) -> str:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[6]  # fully blank: we place our own text boxes
    for title, body in SLIDES:
        slide = prs.slides.add_slide(blank)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
        title_box.text_frame.text = title
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(4.5))
        frame = body_box.text_frame
        for line_no, line in enumerate(body):
            paragraph = frame.paragraphs[0] if line_no == 0 else frame.add_paragraph()
            paragraph.text = line
    prs.save(path)
    return path


if __name__ == "__main__":
    here = os.path.dirname(os.path.abspath(__file__))
    out = os.path.join(here, BASENAME)
    build(out)
    print(f"built {out} ({len(SLIDES)} slides)")
