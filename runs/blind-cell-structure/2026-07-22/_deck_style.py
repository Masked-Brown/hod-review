#!/usr/bin/env python3
"""_deck_style.py — shared visual system + slide builders for the blind-test decks.

One module, one job: give the three blind-test lesson decks (runs/blind-01/) a single,
consistent look and a single set of declarative slide builders, so all three read as the
same designer's hand with only the accent colour changing. The deck scripts import this
module, pass in *content* (including whatever flaws they seed), and never re-implement
layout. This module is style only — it knows nothing about AQA content or seeded flaws.

Design contract enforced here (see the build brief):
  * python-pptx, 16:9 (13.333 x 7.5 in). One accent colour per deck.
  * Full-colour title/divider slides; light-tint content slides; white rounded cards with
    soft shadows; a small spec-code chip; small footer text + slide number.
  * NO full-width header/footer bars, NO edge stripes, NO accent underline beneath titles.
  * Heading font Trebuchet MS, body Calibri. Titles >= 32pt, body 15-18pt.
  * Diagram helper draws internal-labelled nodes + a bounded legend (swatch + text); no
    external leader labels that run off the slide.

Also ships `qa_geometry(prs)` — a deterministic, renderer-free layout check (off-slide
shapes, overlapping cards, estimated text overflow) the deck scripts run before returning.

Stdlib + python-pptx only.
"""

from __future__ import annotations

import math
from dataclasses import dataclass

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------------------
# Constants
# --------------------------------------------------------------------------------------

EMU_PER_INCH = 914400
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
SLIDE_W = int(SLIDE_W_IN * EMU_PER_INCH)
SLIDE_H = int(SLIDE_H_IN * EMU_PER_INCH)

HEAD_FONT = "Trebuchet MS"
BODY_FONT = "Calibri"

MARGIN = 0.62  # outer content margin, inches


# --------------------------------------------------------------------------------------
# Theme
# --------------------------------------------------------------------------------------

@dataclass
class Theme:
    key: str
    accent: RGBColor
    accent_deep: RGBColor
    accent_soft: RGBColor      # very light tint for content-slide backgrounds
    card: RGBColor
    ink: RGBColor
    muted: RGBColor
    line: RGBColor
    spec_code: str
    footer: str


def _rgb(hexstr: str) -> RGBColor:
    return RGBColor.from_string(hexstr)


# One palette per deck. Only the accent family changes; neutrals are shared so the three
# decks are unmistakably one system.
_PALETTES = {
    "A": dict(accent="0E7C66", accent_deep="0A4A3E", accent_soft="EAF4F1"),  # teal
    "B": dict(accent="3F3D8F", accent_deep="27265C", accent_soft="ECECF7"),  # indigo
    "C": dict(accent="B0640F", accent_deep="6F3F09", accent_soft="FBF2E4"),  # amber/ochre
}

_INK = "23272E"
_MUTED = "5C6570"
_LINE = "D8DDE3"
_CARD = "FFFFFF"


def build_theme(key: str, spec_code: str, footer: str) -> Theme:
    """Return the Theme for deck 'A' | 'B' | 'C' with a deck-specific spec code + footer."""
    p = _PALETTES[key]
    return Theme(
        key=key,
        accent=_rgb(p["accent"]),
        accent_deep=_rgb(p["accent_deep"]),
        accent_soft=_rgb(p["accent_soft"]),
        card=_rgb(_CARD),
        ink=_rgb(_INK),
        muted=_rgb(_MUTED),
        line=_rgb(_LINE),
        spec_code=spec_code,
        footer=footer,
    )


# --------------------------------------------------------------------------------------
# Low-level helpers
# --------------------------------------------------------------------------------------

def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    return prs


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_soft_shadow(shape, blur_in=0.075, dist_in=0.045, dir_deg=90,
                    alpha=0.26, color="8A929B") -> None:
    """Inject a soft outer drop shadow into a shape's spPr (python-pptx has no API for it)."""
    spPr = shape._element.spPr
    for el in spPr.findall(qn("a:effectLst")):
        spPr.remove(el)
    blur = int(blur_in * EMU_PER_INCH)
    dist = int(dist_in * EMU_PER_INCH)
    direction = int((dir_deg % 360) * 60000)
    alpha_val = int(alpha * 100000)
    xml = (
        '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:outerShdw blurRad="{blur}" dist="{dist}" dir="{direction}" rotWithShape="0">'
        f'<a:srgbClr val="{color}"><a:alpha val="{alpha_val}"/></a:srgbClr>'
        '</a:outerShdw></a:effectLst>'
    )
    spPr.append(parse_xml(xml))


def _no_shadow(shape) -> None:
    shape.shadow.inherit = False


def rounded_card(slide, x, y, w, h, theme: Theme, fill: RGBColor = None,
                 line: RGBColor = None, line_w=0.75, radius=0.055, shadow=True,
                 name="card::") -> object:
    """A rounded rectangle card. Shapes named 'card::*' are checked for overlap by qa."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill if fill is not None else theme.card
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    if shadow:
        add_soft_shadow(shape)
    else:
        _no_shadow(shape)
    return shape


def plain_rect(slide, x, y, w, h, fill: RGBColor, line: RGBColor = None,
               line_w=0.75, name="deco::") -> object:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    _no_shadow(shape)
    return shape


def _apply_run(run, text, size, color, font=BODY_FONT, bold=False, italic=False):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_par(tf, first, text, size, color, font=BODY_FONT, bold=False, italic=False,
            align=PP_ALIGN.LEFT, space_before=0, space_after=4, level=0,
            line_spacing=1.05, bullet=None, bullet_color=None):
    """Add (or fill) a paragraph. `bullet` (a glyph like '▪') prepends a coloured marker."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.level = level
    if space_before is not None:
        p.space_before = Pt(space_before)
    if space_after is not None:
        p.space_after = Pt(space_after)
    if line_spacing is not None:
        p.line_spacing = line_spacing
    if bullet:
        r0 = p.add_run()
        _apply_run(r0, bullet + "  ", size, bullet_color or color, font=font, bold=True)
    r = p.add_run()
    _apply_run(r, text, size, color, font=font, bold=bold, italic=italic)
    return p


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, name="txt::",
            ml=0.08, mr=0.08, mt=0.04, mb=0.04, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.name = name
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(ml)
    tf.margin_right = Inches(mr)
    tf.margin_top = Inches(mt)
    tf.margin_bottom = Inches(mb)
    return tb, tf


def pill(slide, x, y, w, h, text, fill: RGBColor, text_color: RGBColor,
         size=10.5, line: RGBColor = None, bold=True, font=HEAD_FONT, name="pill::"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    try:
        shape.adjustments[0] = 0.5
    except Exception:
        pass
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.0)
    _no_shadow(shape)
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.09)
    tf.margin_right = Inches(0.09)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_par(tf, True, text, size, text_color, font=font, bold=bold,
            align=PP_ALIGN.CENTER, space_after=0, line_spacing=1.0)
    return shape


# --------------------------------------------------------------------------------------
# Chrome: spec chip, footer, title band
# --------------------------------------------------------------------------------------

def _spec_chip(slide, theme: Theme):
    w = 0.30 + 0.105 * len(theme.spec_code)
    pill(slide, SLIDE_W_IN - MARGIN - w, 0.46, w, 0.34, theme.spec_code,
         fill=theme.accent, text_color=_rgb("FFFFFF"), size=10.5, name="pill::spec")


def _footer(slide, theme: Theme, number: int):
    tb, tf = textbox(slide, MARGIN, SLIDE_H_IN - 0.46, 9.5, 0.32,
                     anchor=MSO_ANCHOR.MIDDLE, name="txt::footer")
    add_par(tf, True, theme.footer, 8.5, theme.muted, font=BODY_FONT, space_after=0)
    tb2, tf2 = textbox(slide, SLIDE_W_IN - MARGIN - 1.2, SLIDE_H_IN - 0.46, 1.2, 0.32,
                       anchor=MSO_ANCHOR.MIDDLE, name="txt::pageno")
    add_par(tf2, True, str(number), 8.5, theme.muted, font=BODY_FONT,
            align=PP_ALIGN.RIGHT, space_after=0)


def content_slide(prs, theme: Theme, title: str, kicker: str = None):
    """Start a light-tint content slide: background, title (no underline), spec chip, footer.

    Returns the slide. Title top edge ~0.46in; content region begins ~1.55in down.
    """
    slide = _blank(prs)
    _set_bg(slide, theme.accent_soft)
    _spec_chip(slide, theme)
    top = 0.46
    if kicker:
        tbk, tfk = textbox(slide, MARGIN, top, 9.0, 0.32, name="txt::kicker")
        add_par(tfk, True, kicker.upper(), 11, theme.accent, font=HEAD_FONT, bold=True,
                space_after=0)
        top += 0.36
    tb, tf = textbox(slide, MARGIN, top, SLIDE_W_IN - 2 * MARGIN - 1.9, 0.85,
                     anchor=MSO_ANCHOR.TOP, name="txt::title")
    add_par(tf, True, title, 32, theme.accent_deep, font=HEAD_FONT, bold=True,
            space_after=0, line_spacing=1.0)
    _footer(slide, theme, len(prs.slides))
    return slide


# --------------------------------------------------------------------------------------
# Slide builders
# --------------------------------------------------------------------------------------

def title_slide(prs, theme: Theme, title, subtitle, meta_lines=None):
    slide = _blank(prs)
    _set_bg(slide, theme.accent)
    # subtle tonal ring motif (bounded — no edge stripe, no bar)
    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.1), Inches(-1.4),
                                  Inches(5.6), Inches(5.6))
    ring.name = "deco::bleed_ring"
    ring.fill.background()
    ring.line.color.rgb = theme.accent_deep
    ring.line.width = Pt(2.0)
    _no_shadow(ring)
    ring2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.35), Inches(-0.15),
                                   Inches(3.1), Inches(3.1))
    ring2.name = "deco::bleed_ring2"
    ring2.fill.background()
    ring2.line.color.rgb = theme.accent_deep
    ring2.line.width = Pt(1.5)
    _no_shadow(ring2)

    pill(slide, MARGIN, 1.55, 2.15, 0.42, theme.spec_code, fill=None,
         text_color=_rgb("FFFFFF"), size=11.5, line=_rgb("FFFFFF"), name="pill::spec")

    tb, tf = textbox(slide, MARGIN, 2.35, 10.6, 2.7, anchor=MSO_ANCHOR.TOP,
                     name="txt::title")
    add_par(tf, True, title, 46, _rgb("FFFFFF"), font=HEAD_FONT, bold=True,
            space_after=6, line_spacing=1.0)
    add_par(tf, False, subtitle, 20, _rgb("EAF1EF"), font=BODY_FONT,
            space_before=2, space_after=0, line_spacing=1.05)

    if meta_lines:
        tbm, tfm = textbox(slide, MARGIN, 5.9, 11.0, 1.0, name="txt::meta")
        for i, line in enumerate(meta_lines):
            add_par(tfm, i == 0, line, 13, _rgb("D8E6E2"), font=BODY_FONT, space_after=2)
    return slide


def divider_slide(prs, theme: Theme, title, subtitle=None):
    slide = _blank(prs)
    _set_bg(slide, theme.accent_deep)
    bar = plain_rect(slide, MARGIN, 3.05, 0.55, 0.10, _rgb("FFFFFF"), name="deco::tick")
    tb, tf = textbox(slide, MARGIN, 3.3, 11.0, 1.8, name="txt::title")
    add_par(tf, True, title, 38, _rgb("FFFFFF"), font=HEAD_FONT, bold=True, space_after=4)
    if subtitle:
        add_par(tf, False, subtitle, 16, _rgb("D8E6E2"), font=BODY_FONT, space_before=2)
    return slide


def staff_slide(prs, theme: Theme, title, rows, spec_points=None, note=None):
    """Two-card 'lesson on a page': left = staff/logistics rows, right = spec this serves."""
    slide = content_slide(prs, theme, title)
    top = 1.7
    lw = 6.35
    left = rounded_card(slide, MARGIN, top, lw, 4.7, theme, name="card::staff")
    tb, tf = textbox(slide, MARGIN + 0.3, top + 0.28, lw - 0.6, 4.2, name="txt::stafftbl")
    for i, (label, value) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(9)
        p.line_spacing = 1.0
        r1 = p.add_run()
        _apply_run(r1, f"{label}   ", 12.5, theme.muted, font=HEAD_FONT, bold=True)
        r2 = p.add_run()
        _apply_run(r2, value, 14, theme.ink, font=BODY_FONT)

    rx = MARGIN + lw + 0.35
    rw = SLIDE_W_IN - MARGIN - rx
    right = rounded_card(slide, rx, top, rw, 4.7, theme, fill=theme.accent_soft,
                         line=theme.line, name="card::spec")
    tb2, tf2 = textbox(slide, rx + 0.28, top + 0.26, rw - 0.56, 4.2, name="txt::specbody")
    add_par(tf2, True, "What this lesson serves", 13, theme.accent_deep, font=HEAD_FONT,
            bold=True, space_after=8)
    for sp in (spec_points or []):
        add_par(tf2, False, sp, 12.5, theme.ink, font=BODY_FONT, bullet="▪",
                bullet_color=theme.accent, space_after=7, line_spacing=1.03)
    if note:
        add_par(tf2, False, note, 11, theme.muted, font=BODY_FONT, italic=True,
                space_before=6, space_after=0)
    return slide


def donow_slide(prs, theme: Theme, title, columns, instruction=None):
    """Retrieval 'Do Now': 3 columns (e.g. Last lesson / Earlier / Stretch)."""
    slide = content_slide(prs, theme, title, kicker="Retrieval · Do Now")
    top = 1.95
    if instruction:
        tbi, tfi = textbox(slide, MARGIN, top - 0.28, 11.0, 0.3, name="txt::donow-instr")
        add_par(tfi, True, instruction, 12.5, theme.muted, font=BODY_FONT, italic=True,
                space_after=0)
    n = len(columns)
    gap = 0.35
    total = SLIDE_W_IN - 2 * MARGIN
    cw = (total - gap * (n - 1)) / n
    ch = 4.35
    for i, (header, lines) in enumerate(columns):
        x = MARGIN + i * (cw + gap)
        rounded_card(slide, x, top, cw, ch, theme, name=f"card::donow{i}")
        hp = pill(slide, x + 0.22, top + 0.24, cw - 0.44, 0.44, header,
                  fill=theme.accent, text_color=_rgb("FFFFFF"), size=12, name=f"pill::dn{i}")
        tb, tf = textbox(slide, x + 0.28, top + 0.86, cw - 0.56, ch - 1.05,
                         name=f"txt::donow{i}")
        for j, line in enumerate(lines):
            add_par(tf, j == 0, line, 13, theme.ink, font=BODY_FONT, bullet="•",
                    bullet_color=theme.accent, space_after=9, line_spacing=1.04)
    return slide


def objectives_slide(prs, theme: Theme, tiers, footnote=None):
    """Tiered ALL/MOST/SOME objectives, each mapped to assessment.

    tiers: list of (tier_label, objective_text, assessment_text).
    """
    slide = content_slide(prs, theme, "Learning objectives")
    top = 1.75
    rowh = 1.32
    gap = 0.18
    label_w = 1.5
    obj_w = 6.6
    ax = MARGIN + label_w + obj_w + 0.3
    aw = SLIDE_W_IN - MARGIN - ax
    # column headers
    tbh, tfh = textbox(slide, MARGIN + label_w + 0.3, top - 0.34, obj_w, 0.3,
                       name="txt::objhdr1")
    add_par(tfh, True, "OBJECTIVE", 10.5, theme.muted, font=HEAD_FONT, bold=True,
            space_after=0)
    tbh2, tfh2 = textbox(slide, ax, top - 0.34, aw, 0.3, name="txt::objhdr2")
    add_par(tfh2, True, "HOW IT IS ASSESSED", 10.5, theme.muted, font=HEAD_FONT, bold=True,
            space_after=0)
    for i, (label, obj, assess) in enumerate(tiers):
        y = top + i * (rowh + gap)
        pill(slide, MARGIN, y + 0.16, label_w, 0.5, label, fill=theme.accent,
             text_color=_rgb("FFFFFF"), size=12.5, name=f"pill::tier{i}")
        rounded_card(slide, MARGIN + label_w + 0.3, y, obj_w, rowh, theme,
                     name=f"card::obj{i}")
        tb, tf = textbox(slide, MARGIN + label_w + 0.5, y + 0.14, obj_w - 0.4, rowh - 0.28,
                         anchor=MSO_ANCHOR.MIDDLE, name=f"txt::obj{i}")
        add_par(tf, True, obj, 13.5, theme.ink, font=BODY_FONT, space_after=0,
                line_spacing=1.03)
        rounded_card(slide, ax, y, aw, rowh, theme, fill=theme.accent_soft,
                     line=theme.line, name=f"card::assess{i}")
        tb2, tf2 = textbox(slide, ax + 0.2, y + 0.14, aw - 0.4, rowh - 0.28,
                           anchor=MSO_ANCHOR.MIDDLE, name=f"txt::assess{i}")
        add_par(tf2, True, assess, 12, theme.accent_deep, font=BODY_FONT, space_after=0,
                line_spacing=1.02)
    if footnote:
        tbf, tff = textbox(slide, MARGIN, top + 3 * (rowh + gap) - 0.02, 11.5, 0.4,
                           name="txt::objnote")
        add_par(tff, True, footnote, 11, theme.muted, font=BODY_FONT, italic=True,
                space_after=0)
    return slide


def keywords_slide(prs, theme: Theme, title, pairs, intro=None):
    """Key words as a two-column grid of term/definition cards."""
    slide = content_slide(prs, theme, title, kicker="Tier-3 vocabulary")
    top = 1.95
    if intro:
        tbi, tfi = textbox(slide, MARGIN, top - 0.3, 11.5, 0.3, name="txt::kw-intro")
        add_par(tfi, True, intro, 12, theme.muted, font=BODY_FONT, italic=True,
                space_after=0)
    cols = 2
    gap_x = 0.35
    gap_y = 0.24
    total = SLIDE_W_IN - 2 * MARGIN
    cw = (total - gap_x) / cols
    rows = math.ceil(len(pairs) / cols)
    ch = min(1.02, (4.5 - gap_y * (rows - 1)) / rows)
    for idx, (term, definition) in enumerate(pairs):
        r = idx // cols
        c = idx % cols
        x = MARGIN + c * (cw + gap_x)
        y = top + r * (ch + gap_y)
        rounded_card(slide, x, y, cw, ch, theme, name=f"card::kw{idx}")
        tb, tf = textbox(slide, x + 0.24, y + 0.13, cw - 0.48, ch - 0.24,
                         anchor=MSO_ANCHOR.MIDDLE, name=f"txt::kw{idx}")
        add_par(tf, True, term, 13.5, theme.accent_deep, font=HEAD_FONT, bold=True,
                space_after=3, line_spacing=1.0)
        add_par(tf, False, definition, 11.5, theme.ink, font=BODY_FONT, space_after=0,
                line_spacing=1.02)
    return slide


def content_slide_bullets(prs, theme: Theme, title, intro=None, bullets=None,
                          side_title=None, side_bullets=None, kicker=None, callout=None):
    """General content: a main card of bullets, an optional narrower side card, an optional
    tinted callout strip along the bottom."""
    slide = content_slide(prs, theme, title, kicker=kicker)
    top = 1.75
    body_h = 4.15 if callout else 4.7
    if side_bullets:
        mw = 7.0
    else:
        mw = SLIDE_W_IN - 2 * MARGIN
    main = rounded_card(slide, MARGIN, top, mw, body_h, theme, name="card::main")
    tb, tf = textbox(slide, MARGIN + 0.32, top + 0.26, mw - 0.64, body_h - 0.5,
                     name="txt::main")
    first = True
    if intro:
        add_par(tf, True, intro, 15.5, theme.ink, font=BODY_FONT, bold=False,
                space_after=10, line_spacing=1.06)
        first = False
    for b in (bullets or []):
        add_par(tf, first, b, 15, theme.ink, font=BODY_FONT, bullet="▪",
                bullet_color=theme.accent, space_after=9, line_spacing=1.06)
        first = False

    if side_bullets:
        sx = MARGIN + mw + 0.35
        sw = SLIDE_W_IN - MARGIN - sx
        rounded_card(slide, sx, top, sw, body_h, theme, fill=theme.accent_soft,
                     line=theme.line, name="card::side")
        tb2, tf2 = textbox(slide, sx + 0.26, top + 0.24, sw - 0.52, body_h - 0.48,
                           name="txt::side")
        if side_title:
            add_par(tf2, True, side_title, 13, theme.accent_deep, font=HEAD_FONT,
                    bold=True, space_after=8)
            sfirst = False
        else:
            sfirst = True
        for b in side_bullets:
            add_par(tf2, sfirst, b, 12.5, theme.ink, font=BODY_FONT, bullet="•",
                    bullet_color=theme.accent, space_after=8, line_spacing=1.04)
            sfirst = False

    if callout:
        cy = top + body_h + 0.22
        rounded_card(slide, MARGIN, cy, SLIDE_W_IN - 2 * MARGIN, 0.68, theme,
                     fill=theme.accent, name="card::callout", shadow=True)
        tbc, tfc = textbox(slide, MARGIN + 0.3, cy + 0.06, SLIDE_W_IN - 2 * MARGIN - 0.6,
                           0.56, anchor=MSO_ANCHOR.MIDDLE, name="txt::callout")
        add_par(tfc, True, callout, 13, _rgb("FFFFFF"), font=BODY_FONT, bold=True,
                space_after=0, line_spacing=1.02)
    return slide


def table_slide(prs, theme: Theme, title, headers, rows, caption=None, kicker=None,
                col_widths=None):
    """A real pptx table with themed header row. Keep cell text short."""
    slide = content_slide(prs, theme, title, kicker=kicker)
    top = 1.95
    nrows = len(rows) + 1
    ncols = len(headers)
    tw = SLIDE_W_IN - 2 * MARGIN
    th = min(0.62 * nrows + 0.15, 4.4)
    gtable = slide.shapes.add_table(nrows, ncols, Inches(MARGIN), Inches(top),
                                    Inches(tw), Inches(th))
    gtable.name = "tbl::main"
    table = gtable.table
    table.first_row = False
    table.horz_banding = False
    if col_widths:
        for ci, cwf in enumerate(col_widths):
            table.columns[ci].width = Inches(tw * cwf)
    # header
    for ci, htext in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme.accent
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.14)
        cell.margin_right = Inches(0.1)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
        tf = cell.text_frame
        tf.word_wrap = True
        add_par(tf, True, htext, 12.5, _rgb("FFFFFF"), font=HEAD_FONT, bold=True,
                space_after=0, line_spacing=1.0)
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = theme.card if ri % 2 else theme.accent_soft
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.14)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            tf = cell.text_frame
            tf.word_wrap = True
            bold = ci == 0
            add_par(tf, True, val, 12, theme.ink, font=BODY_FONT, bold=bold,
                    space_after=0, line_spacing=1.0)
    if caption:
        tbc, tfc = textbox(slide, MARGIN, top + th + 0.16, tw, 0.5, name="txt::tblcap")
        add_par(tfc, True, caption, 11.5, theme.muted, font=BODY_FONT, italic=True,
                space_after=0, line_spacing=1.05)
    return slide


def diagram_node(slide, theme: Theme, x, y, w, h, label, fill: RGBColor,
                 text_color: RGBColor = None, size=11, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
                 name="diag::node"):
    """A shape with an INTERNAL centred label (no external leader lines)."""
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.name = name
    try:
        s.adjustments[0] = 0.18
    except Exception:
        pass
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = _rgb("FFFFFF")
    s.line.width = Pt(1.0)
    _no_shadow(s)
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_par(tf, True, label, size, text_color or _rgb("FFFFFF"), font=HEAD_FONT,
            bold=True, align=PP_ALIGN.CENTER, space_after=0, line_spacing=0.98)
    return s


def draw_legend(slide, theme: Theme, x, y, w, entries, title="Key", row_h=0.42):
    """A bounded legend card: title + rows of [swatch square + text]. entries: (color, text)."""
    h = 0.5 + row_h * len(entries) + 0.16
    rounded_card(slide, x, y, w, h, theme, fill=theme.card, line=theme.line,
                 name="card::legend")
    tb, tf = textbox(slide, x + 0.2, y + 0.13, w - 0.4, 0.32, name="txt::legendttl")
    add_par(tf, True, title, 12, theme.accent_deep, font=HEAD_FONT, bold=True,
            space_after=0)
    for i, (color, text) in enumerate(entries):
        ry = y + 0.52 + i * row_h
        plain_rect(slide, x + 0.22, ry + 0.04, 0.24, 0.24, color, line=theme.line,
                   line_w=0.75, name="deco::swatch")
        tbx, tfx = textbox(slide, x + 0.56, ry - 0.04, w - 0.74, row_h,
                           anchor=MSO_ANCHOR.MIDDLE, name=f"txt::legend{i}")
        add_par(tfx, True, text, 11, theme.ink, font=BODY_FONT, space_after=0,
                line_spacing=1.0)
    return h


def diagram_slide(prs, theme: Theme, title, draw_fn, legend_entries, kicker=None,
                  legend_title="Key", caption=None):
    """A diagram slide. draw_fn(slide, theme, area) draws inside area=(x,y,w,h) inches;
    a bounded legend is placed to the right."""
    slide = content_slide(prs, theme, title, kicker=kicker)
    top = 1.8
    legend_w = 3.5
    diag_w = SLIDE_W_IN - 2 * MARGIN - legend_w - 0.35
    area = (MARGIN, top, diag_w, 4.7)
    # framing card behind the diagram
    rounded_card(slide, area[0], area[1], area[2], area[3], theme, fill=theme.card,
                 line=theme.line, name="card::diagram")
    draw_fn(slide, theme, area)
    lx = MARGIN + diag_w + 0.35
    draw_legend(slide, theme, lx, top, legend_w, legend_entries, title=legend_title)
    if caption:
        tbc, tfc = textbox(slide, lx, top + 3.4, legend_w, 1.2, name="txt::diagcap")
        add_par(tfc, True, caption, 11, theme.muted, font=BODY_FONT, italic=True,
                space_after=0, line_spacing=1.05)
    return slide


def worked_example_slide(prs, theme: Theme, title, i_do, we_do, note=None, kicker=None):
    """I-do / We-do two-card worked example. i_do / we_do are lists of step strings."""
    slide = content_slide(prs, theme, title, kicker=kicker or "Worked example")
    top = 1.8
    gap = 0.35
    cw = (SLIDE_W_IN - 2 * MARGIN - gap) / 2
    ch = 4.35 if note else 4.7
    for i, (head, steps, fillcard) in enumerate([
            ("I do  ·  I model it", i_do, theme.card),
            ("We do  ·  we try it together", we_do, theme.accent_soft)]):
        x = MARGIN + i * (cw + gap)
        rounded_card(slide, x, top, cw, ch, theme, fill=fillcard, line=theme.line,
                     name=f"card::we{i}")
        pill(slide, x + 0.24, top + 0.24, cw - 0.48, 0.44, head, fill=theme.accent,
             text_color=_rgb("FFFFFF"), size=12, name=f"pill::we{i}")
        tb, tf = textbox(slide, x + 0.3, top + 0.86, cw - 0.6, ch - 1.05,
                         name=f"txt::we{i}")
        for j, step in enumerate(steps):
            add_par(tf, j == 0, step, 13, theme.ink, font=BODY_FONT, bullet="›",
                    bullet_color=theme.accent, space_after=9, line_spacing=1.05)
    if note:
        tbn, tfn = textbox(slide, MARGIN, top + ch + 0.16, SLIDE_W_IN - 2 * MARGIN, 0.4,
                           name="txt::wenote")
        add_par(tfn, True, note, 11.5, theme.muted, font=BODY_FONT, italic=True,
                space_after=0)
    return slide


def exam_question_slide(prs, theme: Theme, title, stem_lines, marks, ao,
                        command=None, instruction=None, kicker=None):
    """You-do exam-style question. Marks + AO tag pills, a stem card, an instruction strip."""
    slide = content_slide(prs, theme, title, kicker=kicker or "Independent practice · You do")
    top = 1.85
    ch = 3.7
    rounded_card(slide, MARGIN, top, SLIDE_W_IN - 2 * MARGIN, ch, theme, name="card::exam")
    # tag pills top-right of the card
    tags = []
    if command:
        tags.append(command)
    tags.append(f"{marks}")
    tags.append(ao)
    tx = SLIDE_W_IN - MARGIN - 0.3
    for t in reversed(tags):
        w = 0.34 + 0.1 * len(t)
        tx -= w
        pill(slide, tx, top + 0.26, w, 0.36, t, fill=theme.accent_soft,
             text_color=theme.accent_deep, size=10.5, line=theme.line, name="pill::tag")
        tx -= 0.12
    tb, tf = textbox(slide, MARGIN + 0.34, top + 0.85, SLIDE_W_IN - 2 * MARGIN - 0.68,
                     ch - 1.1, name="txt::examstem")
    for j, line in enumerate(stem_lines):
        add_par(tf, j == 0, line, 15, theme.ink, font=BODY_FONT, space_after=9,
                line_spacing=1.08)
    if instruction:
        iy = top + ch + 0.22
        rounded_card(slide, MARGIN, iy, SLIDE_W_IN - 2 * MARGIN, 0.66, theme,
                     fill=theme.accent, name="card::examinstr")
        tbi, tfi = textbox(slide, MARGIN + 0.3, iy + 0.05, SLIDE_W_IN - 2 * MARGIN - 0.6,
                           0.56, anchor=MSO_ANCHOR.MIDDLE, name="txt::examinstr")
        add_par(tfi, True, instruction, 13, _rgb("FFFFFF"), font=BODY_FONT, bold=True,
                space_after=0, line_spacing=1.02)
    return slide


def plenary_slide(prs, theme: Theme, title, question, options, confidence=True,
                  note=None, kicker=None):
    """Plenary: a hinge MCQ (options = list of (letter, text)) + a confidence RAG strip."""
    slide = content_slide(prs, theme, title, kicker=kicker or "Plenary · Hinge question")
    top = 1.8
    qh = 0.95
    rounded_card(slide, MARGIN, top, SLIDE_W_IN - 2 * MARGIN, qh, theme,
                 fill=theme.accent, name="card::hingeq")
    tbq, tfq = textbox(slide, MARGIN + 0.32, top + 0.1, SLIDE_W_IN - 2 * MARGIN - 0.64,
                       qh - 0.2, anchor=MSO_ANCHOR.MIDDLE, name="txt::hingeq")
    add_par(tfq, True, question, 15.5, _rgb("FFFFFF"), font=HEAD_FONT, bold=True,
            space_after=0, line_spacing=1.03)
    # options 2x2
    oy = top + qh + 0.25
    gap = 0.3
    cw = (SLIDE_W_IN - 2 * MARGIN - gap) / 2
    oh = 0.86
    for i, (letter, otext) in enumerate(options):
        r = i // 2
        c = i % 2
        x = MARGIN + c * (cw + gap)
        y = oy + r * (oh + 0.22)
        rounded_card(slide, x, y, cw, oh, theme, name=f"card::opt{i}")
        pill(slide, x + 0.2, y + 0.21, 0.44, 0.44, letter, fill=theme.accent_soft,
             text_color=theme.accent_deep, size=13, name=f"pill::opt{i}")
        tb, tf = textbox(slide, x + 0.78, y + 0.1, cw - 0.95, oh - 0.2,
                         anchor=MSO_ANCHOR.MIDDLE, name=f"txt::opt{i}")
        add_par(tf, True, otext, 12.5, theme.ink, font=BODY_FONT, space_after=0,
                line_spacing=1.02)
    if confidence:
        rows = (len(options) + 1) // 2
        cy = oy + rows * (oh + 0.22) + 0.04
        tbc, tfc = textbox(slide, MARGIN, cy, SLIDE_W_IN - 2 * MARGIN, 0.4,
                           anchor=MSO_ANCHOR.MIDDLE, name="txt::conf")
        p = tfc.paragraphs[0]
        p.line_spacing = 1.0
        r0 = p.add_run()
        _apply_run(r0, "Confidence:  ", 12, theme.muted, font=HEAD_FONT, bold=True)
        for lab, col in [("Red — reteach me", "C0392B"), ("  Amber — nearly", "C28B0E"),
                         ("  Green — secure", "2E8B57")]:
            r = p.add_run()
            _apply_run(r, "● ", 12, _rgb(col), font=BODY_FONT, bold=True)
            r2 = p.add_run()
            _apply_run(r2, lab.strip() + "   ", 12, theme.ink, font=BODY_FONT)
    return slide


# --------------------------------------------------------------------------------------
# Deterministic, renderer-free geometry QA
# --------------------------------------------------------------------------------------

def _text_overflow(shape, default_pt=14.0):
    """Rough estimate: does the text need more height than the frame gives it?
    Returns a message string if it likely overflows, else None."""
    tf = shape.text_frame
    if not (tf.text or "").strip():
        return None  # empty (e.g. decorative) frame — nothing to overflow
    fw_in = (shape.width or 0) / EMU_PER_INCH
    fh_in = (shape.height or 0) / EMU_PER_INCH
    ml = _emu_to_in(tf.margin_left, 0.1)
    mr = _emu_to_in(tf.margin_right, 0.1)
    mt = _emu_to_in(tf.margin_top, 0.05)
    mb = _emu_to_in(tf.margin_bottom, 0.05)
    avail_w_pt = max(6.0, (fw_in - ml - mr) * 72.0)
    avail_h_pt = max(6.0, (fh_in - mt - mb) * 72.0)
    total_h = 0.0
    for p in tf.paragraphs:
        text = "".join(r.text for r in p.runs) or ""
        sizes = [r.font.size.pt for r in p.runs if r.font.size is not None]
        fs = max(sizes) if sizes else default_pt
        ls = p.line_spacing if isinstance(p.line_spacing, (int, float)) else 1.05
        sa = p.space_after.pt if p.space_after is not None else 3.0
        sb = p.space_before.pt if p.space_before is not None else 0.0
        # average glyph advance ~0.50*fs for Calibri/Trebuchet mixed case
        cpl = max(1, int(avail_w_pt / (0.50 * fs)))
        nlines = max(1, math.ceil(max(1, len(text)) / cpl))
        total_h += sb + nlines * fs * ls + sa
    if total_h > avail_h_pt + 2.0:  # 2pt slack
        return f"~{total_h:.0f}pt needed vs {avail_h_pt:.0f}pt available"
    return None


def _emu_to_in(v, default):
    if v is None:
        return default
    try:
        return v / EMU_PER_INCH
    except Exception:
        return default


def _rects_overlap(a, b):
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ix = max(0, min(ax + aw, bx + bw) - max(ax, bx))
    iy = max(0, min(ay + ah, by + bh) - max(ay, by))
    return ix * iy


def qa_geometry(prs, default_font_pt=14.0):
    """Return a list of layout problems (empty == clean). Checks, per slide:
      * shapes off the slide canvas,
      * estimated text overflow of text frames,
      * overlap between sibling 'card::' shapes.
    """
    problems = []
    tol = int(0.03 * EMU_PER_INCH)
    for si, slide in enumerate(prs.slides, 1):
        cards = []
        for sh in slide.shapes:
            if sh.left is None or sh.top is None:
                continue
            l, t = sh.left, sh.top
            w = sh.width or 0
            h = sh.height or 0
            bleed = sh.name.startswith("deco::bleed")
            if not bleed and (
                    l < -tol or t < -tol or (l + w) > SLIDE_W + tol or (t + h) > SLIDE_H + tol):
                problems.append(
                    f"S{si}: '{sh.name}' off-slide "
                    f"(l={l/EMU_PER_INCH:.2f} t={t/EMU_PER_INCH:.2f} "
                    f"r={(l+w)/EMU_PER_INCH:.2f} b={(t+h)/EMU_PER_INCH:.2f})")
            if getattr(sh, "has_text_frame", False):
                msg = _text_overflow(sh, default_font_pt)
                if msg:
                    problems.append(f"S{si}: '{sh.name}' text overflow ({msg})")
            if sh.name.startswith("card::"):
                cards.append((sh.name, (l, t, w, h)))
        for i in range(len(cards)):
            for j in range(i + 1, len(cards)):
                if _rects_overlap(cards[i][1], cards[j][1]) > int(0.02 * EMU_PER_INCH) ** 2:
                    problems.append(
                        f"S{si}: cards '{cards[i][0]}' and '{cards[j][0]}' overlap")
    return problems


if __name__ == "__main__":  # pragma: no cover - smoke build
    prs = new_deck()
    th = build_theme("A", "AQA 3.2.1", "smoke test")
    title_slide(prs, th, "Smoke test", "the shared style module", ["one", "two"])
    content_slide_bullets(prs, th, "A content slide", intro="Intro line.",
                          bullets=["First point here.", "Second point here."])
    import os
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "_smoke.pptx")
    prs.save(out)
    print("qa:", qa_geometry(prs))
    print("wrote", out)
